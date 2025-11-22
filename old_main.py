import os
import re
import json
import yaml
import asyncio
import argparse
import logging
from typing import Dict, List, Any, Optional, Tuple
from pprint import pprint
from pathlib import Path
import openai

# 尝试导入config模块，如果不存在则创建默认配置
try:
    from config import example_format, config
except ImportError:
    # 创建默认配置
    example_format = ""
    config = {
        "threshold": 0.8,
        "openai_api_key": os.getenv("OPENAI_API_KEY", "")
    }

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# 初始化OpenAI客户端
openai_client = openai.OpenAI(api_key=config.get("openai_api_key", os.getenv("OPENAI_API_KEY", "")))

# 导入mineru相关模块
from mineru.cli.common import do_parse, read_fn
from mineru.utils.enum_class import MakeMode

class AI_Profiler:
    def __init__(self, knowledge_dir="knowledges", output_dir="outputs", mid_results_dir="mid_results", config=None):
        self.knowledge_dir = knowledge_dir
        self.output_dir = output_dir
        self.mid_results_dir = mid_results_dir
        
        self.config = config or {}
        self.supported_formats = ['.pptx', '.pdf']
        
        for dir_path in [output_dir, mid_results_dir]:
            Path(dir_path).mkdir(exist_ok=True)
    
    def extract_text_from_ppt(self, file_path: str) -> str:
        """
        使用mineru工具从PPT文件中提取文本并转换为markdown格式
        
        Args:
            file_path: PPT文件路径
            
        Returns:
            str: 提取的markdown文本内容
        """
        try:
            # 创建临时输出目录
            output_dir = os.path.join(os.path.dirname(file_path), 'mineru_output')
            os.makedirs(output_dir, exist_ok=True)
            
            # 获取文件名（不含扩展名）
            file_name = os.path.splitext(os.path.basename(file_path))[0]
            
            # 读取文件字节
            file_bytes = read_fn(file_path)
            
            # 使用mineru进行文档解析
            do_parse(
                output_dir=output_dir,
                pdf_file_names=[file_name],
                pdf_bytes_list=[file_bytes],
                p_lang_list=['ch'],  # 假设中文文档
                backend="pipeline",
                parse_method="auto",
                formula_enable=True,
                table_enable=True,
                f_draw_layout_bbox=False,  # 不需要可视化
                f_draw_span_bbox=False,
                f_dump_md=True,
                f_dump_orig_pdf=False,
                f_dump_content_list=False,
                f_dump_middle_json=False,
                f_dump_model_output=False,
                f_make_md_mode=MakeMode.MM_MD
            )
            
            # 读取生成的markdown文件
            md_file_path = os.path.join(output_dir, file_name, 'auto', f'{file_name}.md')
            with open(md_file_path, 'r', encoding='utf-8') as f:
                markdown_content = f.read()
            
            return markdown_content
            
        except Exception as e:
            logger.error(f"Error extracting text from PPT using mineru: {e}")
            # 如果mineru失败，尝试使用原始的python-pptx作为备份
            try:
                from pptx import Presentation
                
                prs = Presentation(file_path)
                text_runs = []
                
                for slide_idx, slide in enumerate(prs.slides, 1):
                    text_runs.append(f"## Slide {slide_idx}")
                    
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            # 处理段落格式
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    # 判断标题级别
                                    level = paragraph.level
                                    if level == 0 and paragraph.font.size and paragraph.font.size.pt > 24:
                                        text_runs.append(f"### {paragraph.text.strip()}")
                                    elif level == 0:
                                        text_runs.append(f"#### {paragraph.text.strip()}")
                                    else:
                                        text_runs.append(f"{'  ' * level}- {paragraph.text.strip()}")
                    
                    text_runs.append("")  # 添加空行分隔幻灯片
                
                return "\n".join(text_runs)
                
            except Exception as fallback_error:
                logger.error(f"Fallback method also failed: {fallback_error}")
                return ""
    
    def clean_and_structure_content(self, content: str) -> str:
        """清理和结构化提取的内容"""
        content = re.sub(r'\n\s*\n\s*\n', '\n\n', content)
        content = content.strip()
        
        # 添加标题
        if not content.startswith('#'):
            content = f"# AI Analysis Report\n\n{content}"
        
        return content
    
    def generate_skills(self, structured_content: Dict) -> Dict:
        content_summary = "\n\n".join([
            f"第{page['page_number']}页:\n{page['content'][:200]}..." if len(page['content']) > 200 else f"第{page['page_number']}页:\n{page['content']}"
            for page in structured_content[:10]
        ])
        
        if len(structured_content) > 10:
            content_summary += f"\n\n... 以及其他 {len(structured_content) - 10} 页内容"
        
        from config import example_format
        
        prompt = f"""
        你是一位教育专家，任务是从课程PPT内容中识别学生需要掌握的技能点，并将这些技能点分组(group)，请着重注意，每个技能点可以出现在多个group中。        
        每个group的目标是可以用一道选择或判断题目联合考察这一group内的所有或部分技能点，每个group必须至少含有3个技能点。
        每个group的类型(type)可以是：
        - Dependency：表示这一个group中的某些技能点是此group中某些其他技能点的前置条件或基础知识，故可以一起考查。
        - Progressive：表示这一个group中的技能点是关于某一知识的不同深度的内容，如“掌握某个知识的定义”、“理解这个知识的应用方法”等技能点就可以出现在同一个Progressive组中。
        - Contrast：表示这一个group中的技能点是关于某些知识的对比理解，如A和B是两个相对的知识点，那么“理解A”、“理解B”等技能点就可以出现在同一个Contrast组中。
        - Application：表示这一个group中的技能点是关于某些知识的实际应用，这些技能点可能是运用不同方法解决相同问题场景，也可能是运用同一种方法解决不同问题。
        - Related：表示这一个group中的技能点之间存在着不属于以上四种关系的相关性，也可以在同一道题目中考察。

        请基于以下PPT内容，提取出：
        1. 所有学生需要掌握的具体技能点列表，每个技能点应该以"能够..."的格式表述
        2. 将相关的技能点分组，每组包含：
           - 组ID（如Group1_Definition）
           - 组类型
           - 组目标（target，描述该组技能点的总体目标）
           - 包含的技能点ID列表
        
        参考格式：
        {example_format}
        
        请确保：
        - 技能点描述清晰、具体，聚焦于学生能够执行的能力
        - 技能点覆盖PPT中的所有核心概念和方法
        - 分组逻辑合理，便于出题时合并考查相关技能点
        - 输出必须是有效的JSON格式
        - 输出中skills应对应一个列表，其中每个元素的格式为：
        "0: 能够形式化定义上下文无关文法 G = (V, T, P, S) 的四个组成部分及其约束条件",
        
        PPT内容：
        ```
        {content_summary}
        ```
        """
        
        try:
            logger.info("开始调用大模型识别技能点...")
            response = openai_client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "你是一位教育评估专家，擅长从教学内容中提取学习目标和技能点。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=2000
            )
            content = response.choices[0].message.content.strip()
            
            if content.startswith('```json') and content.endswith('```'):
                content = content[7:-3].strip()
            elif content.startswith('```') and content.endswith('```'):
                content = content[3:-3].strip()
            
            result = json.loads(content)
            
            mid_file = os.path.join(self.mid_results_dir, "identified_skills.json")
            with open(mid_file, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            logger.info(f"已保存识别的技能点到 {mid_file}")
            return result
        except Exception as e:
            logger.error(f"调用大模型识别技能点时出错: {e}")
    
    def convert_to_yml(self, skills_data: Dict, output_file: str) -> None:
        yml_data = {}
        yml_data["skills"] = []
        for idx, skill in enumerate(skills_data["skills"]):
            if isinstance(skill, dict):
                yml_data["skills"].append(skill)
            else:
                if isinstance(skill, str) and ':' in skill:
                    parts = skill.split(':', 1)
                    if len(parts) == 2:
                        try:
                            skill_idx = int(parts[0].strip())
                            skill_desc = parts[1].strip()
                            yml_data["skills"].append({skill_idx: skill_desc})
                        except ValueError:
                            yml_data["skills"].append({idx: skill})
                else:
                    yml_data["skills"].append({idx: str(skill)})
        
        yml_data["groups"] = skills_data["groups"]
        
        output_path = os.path.join(self.output_dir, output_file)
        with open(output_path, 'w', encoding='utf-8') as f:
            yaml.dump(yml_data, f, allow_unicode=True, default_flow_style=False, sort_keys=False)
        logger.info(f"已生成YML文件: {output_path}")
    
    def validate_yml_structure(self, yml_file: str) -> bool:
        try:
            with open(yml_file, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
            
            if "skills" not in data or "groups" not in data:
                logger.error("YML文件缺少必需的'skills'或'groups'键")
                return False
            
            if not isinstance(data["skills"], list):
                logger.error("skills应为列表格式")
                return False
            
            for i, skill in enumerate(data["skills"]):
                if not isinstance(skill, dict) or len(skill) != 1:
                    logger.error(f"技能点 {i} 格式错误")
                    return False
            
            if not isinstance(data["groups"], dict):
                logger.error("groups应为字典格式")
                return False
            
            for group_id, group in data["groups"].items():
                if not isinstance(group, dict):
                    logger.error(f"群组 {group_id} 格式错误")
                    return False
                if "type" not in group or "target" not in group or "skills" not in group:
                    logger.error(f"群组 {group_id} 缺少必需的属性")
                    return False
                if not isinstance(group["skills"], list):
                    logger.error(f"群组 {group_id} 的skills应为列表格式")
                    return False
            
            logger.info(f"YML文件 {yml_file} 结构验证通过")
            return True
        except Exception as e:
            logger.error(f"验证YML文件结构时出错: {e}")
            return False
    
    def run_workflow(self, input_file: str, output_file: str) -> bool:
        """运行完整的工作流"""
        try:
            # 1. 从PPT提取文本
            logger.info(f"Extracting text from {input_file}...")
            content = self.extract_text_from_ppt(input_file)
            
            if not content:
                logger.error("Failed to extract content from file")
                return False
            
            # 2. 清理和结构化内容
            logger.info("Cleaning and structuring content...")
            structured_content = self.clean_and_structure_content(content)
            
            # 保存结构化内容
            mid_file = os.path.join(self.mid_results_dir, "structured_content.json")
            with open(mid_file, 'w', encoding='utf-8') as f:
                json.dump([{"page_number": i+1, "content": page} for i, page in enumerate(structured_content.split("\n\n## Slide"))], f, ensure_ascii=False, indent=2)
            logger.info(f"已保存结构化内容到 {mid_file}")
            
            # 3. 生成技能点
            logger.info("Generating skills...")
            skills_data = self.generate_skills([{"page_number": i+1, "content": page} for i, page in enumerate(structured_content.split("\n\n## Slide"))])
            
            # 4. 转换为YML格式
            logger.info("Converting to YML format...")
            self.convert_to_yml(skills_data, output_file)
            
            # 5. 验证YML文件结构
            yml_path = os.path.join(self.output_dir, output_file)
            if self.validate_yml_structure(yml_path):
                logger.info("Success")
                return True
            else:
                logger.error("YML文件结构验证失败")
                return False
        except Exception as e:
            logger.error(f"工作流执行失败: {e}")
            return False

def main():
    """主函数"""
    # 创建命令行参数解析器
    parser = argparse.ArgumentParser(description='AI Profiler - 从PPT/PDF文件生成技能分析报告')
    parser.add_argument('-i', '--input', type=str, help='输入文件路径')
    parser.add_argument('-d', '--directory', type=str, help='输入目录路径')
    parser.add_argument('-o', '--output', type=str, help='输出文件路径')
    
    # 解析命令行参数
    args = parser.parse_args()
    
    # 初始化AI分析器
    profiler = AI_Profiler()
    
    # 处理输入文件
    if args.input:
        if not os.path.exists(args.input):
            print(f"Error: File {args.input} does not exist")
            return
        
        # 确定输出文件路径
        if args.output:
            output_path = args.output
        else:
            output_path = os.path.splitext(args.input)[0] + '.yml'
        
        # 运行工作流
        profiler.run_workflow(args.input, output_path)
    
    # 处理目录
    elif args.directory:
        if not os.path.exists(args.directory):
            print(f"Error: Directory {args.directory} does not exist")
            return
        
        # 查找支持的文件
        files = find_files(args.directory, profiler.supported_formats)
        
        if not files:
            print(f"No supported files found in {args.directory}")
            return
        
        print(f"Found {len(files)} files to process:")
        for file in files:
            print(f"  - {file}")
        
        # 为每个文件运行工作流
        for file in files:
            output_path = os.path.splitext(file)[0] + '.yml'
            print(f"\nProcessing {file}...")
            profiler.run_workflow(file, output_path)
    
    else:
        # 尝试在当前目录查找
        current_dir = os.getcwd()
        files = find_files(current_dir, profiler.supported_formats)
        
        if not files:
            print(f"No supported files found in current directory")
            parser.print_help()
            return
        
        print(f"Found {len(files)} files to process:")
        for file in files:
            print(f"  - {file}")
        
        # 为每个文件运行工作流
        for file in files:
            output_path = os.path.splitext(file)[0] + '.yml'
            print(f"\nProcessing {file}...")
            profiler.run_workflow(file, output_path)

def find_files(directory: str, extensions: List[str]) -> List[str]:
    """
    查找目录中具有指定扩展名的文件
    
    Args:
        directory: 要搜索的目录路径
        extensions: 要查找的文件扩展名列表
        
    Returns:
        List[str]: 找到的文件路径列表
    """
    files = []
    for root, _, filenames in os.walk(directory):
        for filename in filenames:
            if any(filename.lower().endswith(ext) for ext in extensions):
                files.append(os.path.join(root, filename))
    return files

def get_coverage_metrics(structured_content_path: str, skills_data: Dict) -> Dict:
    """
    计算技能点覆盖率指标
    
    Args:
        structured_content_path: 结构化内容文件路径
        skills_data: 识别出的技能点数据
        
    Returns:
        Dict: 覆盖率指标，包含coverage（覆盖率）、threshold（阈值）、is_valid（是否有效）
    """
    try:
        # 读取结构化内容
        with open(structured_content_path, 'r', encoding='utf-8') as f:
            structured_content = json.load(f)
        
        # 简单的覆盖率计算（假设每个技能点覆盖一个知识点）
        total_content_sections = len(structured_content)
        identified_skills_count = len(skills_data.get("skills", []))
        
        # 计算覆盖率（这里使用简单的比例，实际可能需要更复杂的算法）
        coverage = min(identified_skills_count / total_content_sections if total_content_sections > 0 else 0, 1.0)
        threshold = config.get("threshold", 0.8)
        
        return {
            "coverage": coverage,
            "threshold": threshold,
            "is_valid": coverage >= threshold
        }
    except Exception as e:
        logger.error(f"计算覆盖率指标时出错: {e}")
        return {
            "coverage": 0.0,
            "threshold": config.get("threshold", 0.8),
            "is_valid": False
        }

if __name__ == "__main__":
    main()