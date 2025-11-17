import os
import re
import json
import yaml
import time
import logging
from pathlib import Path
from typing import List, Dict, Any
import PyPDF2

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("conversion.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

from client import openai_client

class AI_Profiler:
    def __init__(self, knowledge_dir="knowledges", output_dir="outputs", mid_results_dir="mid_results"):
        self.knowledge_dir = knowledge_dir
        self.output_dir = output_dir
        self.mid_results_dir = mid_results_dir
        
        for dir_path in [output_dir, mid_results_dir]:
            Path(dir_path).mkdir(exist_ok=True)
    
    def extract_text_from_ppt(self, file_path: str) -> str:
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return self._extract_from_pdf(file_path)
        elif file_extension == '.pptx':
            try:
                return self._extract_from_pptx(file_path)
            except ImportError:
                logger.warning("未找到python-pptx模块")
        else:
            logger.error(f"不支持的文件格式: {file_extension}")
            return ""
    
    def _extract_from_pdf(self, pdf_path: str) -> str:
        text = ""
        try:
            with open(pdf_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n=== 第{page_num + 1}页 ===\n"
                        text += page_text
            logger.info(f"从PDF提取了 {len(text)} 个字符")
            return text
        except Exception as e:
            logger.error(f"提取PDF内容时出错: {e}")
            return ""
    
    def _extract_from_pptx(self, pptx_path: str) -> str:
        from pptx import Presentation
        
        text = ""
        try:
            prs = Presentation(pptx_path)
            for i, slide in enumerate(prs.slides):
                text += f"\n=== 第{i + 1}页 ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        text += shape.text + "\n"
            logger.info(f"从PPTX提取了 {len(text)} 个字符")
            return text
        except Exception as e:
            logger.error(f"提取PPTX内容时出错: {e}")
            return ""

    def clean_and_structure_content(self, raw_text: str) -> Dict:
        text = re.sub(r'\s+', ' ', raw_text)
        text = re.sub(r'\n\s*\n', '\n\n', text)
        
        pages = re.split(r'=== 第\d+页 ===', text)
        pages = [page.strip() for page in pages if page.strip()]
        
        structured_content = []
        for i, page in enumerate(pages):
            structured_content.append({
                "page_number": i + 1,
                "content": page
            })
        
        mid_file = os.path.join(self.mid_results_dir, "structured_content.json")
        with open(mid_file, 'w', encoding='utf-8') as f:
            json.dump(structured_content, f, ensure_ascii=False, indent=2)
        logger.info(f"已保存结构化内容到 {mid_file}")
        
        return structured_content
    
    def generate_skills(self, structured_content: Dict) -> Dict:
        content_summary = "\n\n".join([
            f"第{page['page_number']}页:\n{page['content'][:200]}..." if len(page['content']) > 200 else f"第{page['page_number']}页:\n{page['content']}"
            for page in structured_content[:10]
        ])
        
        if len(structured_content) > 10:
            content_summary += f"\n\n... 以及其他 {len(structured_content) - 10} 页内容"
        
        from config import example_format
        
        prompt = f"""
        你是一位教育专家，任务是从课程PPT内容中识别学生需要掌握的技能点，并将这些技能点分组(group)，请着重注意，每个技能点可以出现在多个group中。        
        每个group的目标是可以用一道选择或判断题目联合考察这一group内的所有或部分技能点，每个group必须至少含有3个技能点。
        每个group的类型(type)可以是：
        - Dependency：表示这一个group中的某些技能点是此group中某些其他技能点的前置条件或基础知识，故可以一起考查。
        - Progressive：表示这一个group中的技能点是关于某一知识的不同深度的内容，如“掌握某个知识的定义”、“理解这个知识的应用方法”等技能点就可以出现在同一个Progressive组中。
        - Contrast：表示这一个group中的技能点是关于某些知识的对比理解，如A和B是两个相对的知识点，那么“理解A”、“理解B”等技能点就可以出现在同一个Contrast组中。
        - Application：表示这一个group中的技能点是关于某些知识的实际应用，这些技能点可能是运用不同方法解决相同问题场景，也可能是运用同一种方法解决不同问题。
        - Related：表示这一个group中的技能点之间存在着不属于以上四种关系的相关性，也可以在同一道题目中考察。

        请基于以下PPT内容，提取出：
        1. 所有学生需要掌握的具体技能点列表，每个技能点应该以"能够..."的格式表述
        2. 将相关的技能点分组，每组包含：
           - 组ID（如Group1_Definition）
           - 组类型
           - 组目标（target，描述该组技能点的总体目标）
           - 包含的技能点ID列表
        
        参考格式：
        {example_format}
        
        请确保：
        - 技能点描述清晰、具体，聚焦于学生能够执行的能力
        - 技能点覆盖PPT中的所有核心概念和方法
        - 分组逻辑合理，便于出题时合并考查相关技能点
        - 输出必须是有效的JSON格式
        - 输出中skills应对应一个列表，其中每个元素的格式为：
        "0: 能够形式化定义上下文无关文法 G = (V, T, P, S) 的四个组成部分及其约束条件",
        
        PPT内容：
        ```
        {content_summary}
        ```
        """
        
        try:
            logger.info("开始调用大模型识别技能点...")
            response = openai_client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "你是一位教育评估专家，擅长从教学内容中提取学习目标和技能点。"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.3,
                max_tokens=2000
            )
            content = response.choices[0].message.content.strip()
            
            if content.startswith('```json') and content.endswith('```'):
                content = content[7:-3].strip()
            elif content.startswith('```') and content.endswith('```'):
                content = content[3:-3].strip()
            
            result = json.loads(content)
            
            mid_file = os.path.join(self.mid_results_dir, "identified_skills.json")
            with open(mid_file, 'w', encoding='utf-8') as f:
                json.dump(result, f, ensure_ascii=False, indent=2)
            logger.info(f"已保存识别的技能点到 {mid_file}")
            return result
        except Exception as e:
            logger.error(f"调用大模型识别技能点时出错: {e}")
    
    def convert_to_yml(self, skills_data: Dict, output_file: str) -> None:
        yml_data = {}
        yml_data["skills"] = []
        for idx, skill in enumerate(skills_data["skills"]):
            if isinstance(skill, dict):
                yml_data["skills"].append(skill)
            else:
                if isinstance(skill, str) and ':' in skill:
                    parts = skill.split(':', 1)
                    if len(parts) == 2:
                        try:
                            skill_idx = int(parts[0].strip())
                            skill_desc = parts[1].strip()
                            yml_data["skills"].append({skill_idx: skill_desc})
                        except ValueError:
                            yml_data["skills"].append({idx: skill})
                else:
                    yml_data["skills"].append({idx: str(skill)})
        
        yml_data["groups"] = skills_data["groups"]
        
        output_path = os.path.join(self.output_dir, output_file)
        with open(output_path, 'w', encoding='utf-8') as f:
            yaml.dump(yml_data, f, allow_unicode=True, default_flow_style=False, sort_keys=False)
        logger.info(f"已生成YML文件: {output_path}")
    
    def validate_yml_structure(self, yml_file: str) -> bool:
        try:
            with open(yml_file, 'r', encoding='utf-8') as f:
                data = yaml.safe_load(f)
            
            if "skills" not in data or "groups" not in data:
                logger.error("YML文件缺少必需的'skills'或'groups'键")
                return False
            
            if not isinstance(data["skills"], list):
                logger.error("skills应为列表格式")
                return False
            
            for i, skill in enumerate(data["skills"]):
                if not isinstance(skill, dict) or len(skill) != 1:
                    logger.error(f"技能点 {i} 格式错误")
                    return False
            
            if not isinstance(data["groups"], dict):
                logger.error("groups应为字典格式")
                return False
            
            for group_id, group in data["groups"].items():
                if not isinstance(group, dict):
                    logger.error(f"群组 {group_id} 格式错误")
                    return False
                if "type" not in group or "target" not in group or "skills" not in group:
                    logger.error(f"群组 {group_id} 缺少必需的属性")
                    return False
                if not isinstance(group["skills"], list):
                    logger.error(f"群组 {group_id} 的skills应为列表格式")
                    return False
            
            logger.info(f"YML文件 {yml_file} 结构验证通过")
            return True
        except Exception as e:
            logger.error(f"验证YML文件结构时出错: {e}")
            return False
    
    def run_workflow(self, input_file: str, output_file: str) -> bool:
        try:
            logger.info(f"开始处理文件: {input_file}")
            
            # PPT to raw text
            raw_text = self.extract_text_from_ppt(input_file)
            if not raw_text:
                logger.error("无法提取PPT内容")
                return False
            
            # cleaning raw text
            structured_content = self.clean_and_structure_content(raw_text)
            
            # text to skills
            skills_data = self.generate_skills(structured_content)
            
            # skills to groups
            self.convert_to_yml(skills_data, output_file)
            
            # validation
            yml_path = os.path.join(self.output_dir, output_file)
            if self.validate_yml_structure(yml_path):
                logger.info("Success")
                return True
            else:
                logger.error("YML文件结构验证失败")
                return False
        except Exception as e:
            logger.error(f"工作流执行失败: {e}")
            return False

def main():
    client = AI_Profiler()
    
    ppt_files = [f for f in os.listdir(client.knowledge_dir) 
                if f.lower().endswith(('.pptx', '.pdf'))]
    
    if not ppt_files:
        logger.error(f"在 {client.knowledge_dir} 目录中未找到PPT或PDF文件")
        return
    
    input_file = os.path.join(client.knowledge_dir, ppt_files[3])
    output_file = f"{os.path.splitext(os.path.basename(input_file))[0]}_skills.yml"
    
    client.run_workflow(input_file, output_file)

if __name__ == "__main__":
    main()